import os
import math
import random
from PIL import Image, ImageDraw
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE

def generate_background(width=1920, height=1080, filename="background.png"):
    # Generate a small gradient image first
    sw, sh = 192, 108
    grad = Image.new("RGB", (sw, sh))
    c_center = (38, 40, 48)  # Slate Gray
    c_edge = (6, 6, 8)       # Almost Black
    
    pixels = []
    for y in range(sh):
        for x in range(sw):
            dx = (x - sw / 2.0) / (sw / 2.0)
            dy = (y - sh / 2.0) / (sh / 2.0)
            dist = math.sqrt(dx*dx + dy*dy)
            factor = max(0.0, min(1.0, dist))
            
            r = int(c_center[0] + (c_edge[0] - c_center[0]) * factor)
            g = int(c_center[1] + (c_edge[1] - c_center[1]) * factor)
            b = int(c_center[2] + (c_edge[2] - c_center[2]) * factor)
            pixels.append((r, g, b))
            
    grad.putdata(pixels)
    img = grad.resize((width, height), Image.Resampling.BILINEAR)
    
    # Generate fine grain noise
    noise_sz = 512
    noise = Image.new("L", (noise_sz, noise_sz))
    noise_pixels = [128 + random.randint(-8, 8) for _ in range(noise_sz * noise_sz)]
    noise.putdata(noise_pixels)
    
    noise_large = Image.new("L", (width, height))
    for ty in range(0, height, noise_sz):
        for tx in range(0, width, noise_sz):
            noise_large.paste(noise, (tx, ty))
            
    noise_rgb = noise_large.convert("RGB")
    final_img = Image.blend(img, noise_rgb, 0.05)
    final_img.save(filename, "PNG")

def evaluate_bezier(p0, p1, p2, p3, steps=50):
    points = []
    for i in range(steps + 1):
        t = i / float(steps)
        x = (1-t)**3 * p0[0] + 3*(1-t)**2 * t * p1[0] + 3*(1-t) * t**2 * p2[0] + t**3 * p3[0]
        y = (1-t)**3 * p0[1] + 3*(1-t)**2 * t * p1[1] + 3*(1-t) * t**2 * p2[1] + t**3 * p3[1]
        points.append((x, y))
    return points

def evaluate_arc(cx, cy, r, start_ang, end_ang, steps=50):
    points = []
    for i in range(steps + 1):
        t = i / float(steps)
        ang = start_ang + t * (end_ang - start_ang)
        x = cx + r * math.cos(ang)
        y = cy + r * math.sin(ang)
        points.append((x, y))
    return points

def generate_logo(size=1000, filename="logo.png"):
    scale = size / 100.0
    img = Image.new("RGBA", (size, size), (0, 0, 0, 0))
    draw = ImageDraw.Draw(img)
    
    cx, cy = 50.0 * scale, 50.0 * scale
    r = 27.73 * scale
    
    p1_arc = evaluate_arc(cx, cy, r, 2.7022, 5.8348, steps=50)
    p1_bez = evaluate_bezier((75.0*scale, 38.0*scale), (65.0*scale, 37.0*scale), (44.0*scale, 45.0*scale), (25.0*scale, 62.0*scale), steps=50)
    p1 = p1_arc + p1_bez
    draw.polygon(p1, fill=(255, 255, 255, 255))
    
    p2_arc = evaluate_arc(cx, cy, r, -0.4402, 2.7022, steps=50)
    p2_bez = evaluate_bezier((25.0*scale, 62.0*scale), (35.0*scale, 63.0*scale), (56.0*scale, 55.0*scale), (75.0*scale, 38.0*scale), steps=50)
    p2 = p2_arc + p2_bez
    draw.polygon(p2, fill=(255, 255, 255, 255))
    
    p3_fwd = evaluate_bezier((16.0*scale, 66.0*scale), (36.0*scale, 50.0*scale), (58.0*scale, 36.0*scale), (84.0*scale, 26.0*scale), steps=50)
    p3_bwd = evaluate_bezier((84.0*scale, 26.0*scale), (64.0*scale, 46.0*scale), (42.0*scale, 60.0*scale), (16.0*scale, 66.0*scale), steps=50)
    p3 = p3_fwd + p3_bwd
    draw.polygon(p3, fill=(255, 121, 0, 255))
    
    img_resized = img.resize((256, 256), Image.Resampling.LANCZOS)
    img_resized.save(filename, "PNG")

def create_presentation():
    # Generate background and logo images first
    generate_background()
    generate_logo()

    prs = Presentation()
    # Set slide dimensions to widescreen 16:9
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    
    # Theme Colors
    BG_COLOR = RGBColor(10, 10, 12)        # Deep Dark Slate
    TEXT_WHITE = RGBColor(255, 255, 255)   # Main Text White
    TEXT_MUTED = RGBColor(161, 161, 170)   # Muted/Body Text (Zinc-400)
    COLOR_ORANGE = RGBColor(255, 121, 0)   # Suzano IT Orange Accent
    COLOR_INDIGO = COLOR_ORANGE
    COLOR_SKY = COLOR_ORANGE
    COLOR_RED = RGBColor(239, 68, 68)      # Warning Accent
    COLOR_CARD_BG = RGBColor(18, 19, 22)   # Card Background
    
    # Slide layouts (6 is blank in default template)
    blank_layout = prs.slide_layouts[6]
    
    def set_dark_background(slide):
        # Add the background image shape
        slide.shapes.add_picture('background.png', 0, 0, prs.slide_width, prs.slide_height)
        
    def add_header(slide, title_text, category_text="PLATAFORMA SAAS"):
        # Category tag
        cat_box = slide.shapes.add_textbox(Inches(0.8), Inches(0.4), Inches(8), Inches(0.4))
        cat_tf = cat_box.text_frame
        cat_tf.word_wrap = True
        p_cat = cat_tf.paragraphs[0]
        p_cat.text = category_text.upper()
        p_cat.font.name = 'Arial'
        p_cat.font.size = Pt(10)
        p_cat.font.bold = True
        p_cat.font.color.rgb = COLOR_SKY
        
        # Slide Title
        title_box = slide.shapes.add_textbox(Inches(0.8), Inches(0.7), Inches(11.5), Inches(0.8))
        title_tf = title_box.text_frame
        title_tf.word_wrap = True
        title_tf.margin_left = Inches(0)
        p_title = title_tf.paragraphs[0]
        p_title.text = title_text
        p_title.font.name = 'Arial'
        p_title.font.size = Pt(28)
        p_title.font.bold = True
        p_title.font.color.rgb = TEXT_WHITE
        
    def add_footer(slide, current_page, total_pages=15):
        # Footer text
        footer_box = slide.shapes.add_textbox(Inches(0.8), Inches(7.0), Inches(11.7), Inches(0.4))
        footer_tf = footer_box.text_frame
        p_foot = footer_tf.paragraphs[0]
        p_foot.text = f"SUZANO IT   |   Tecnologia que move frotas.                                                                                          Slide {current_page} de {total_pages}"
        p_foot.font.name = 'Arial'
        p_foot.font.size = Pt(8.5)
        p_foot.font.color.rgb = TEXT_MUTED
        
    def draw_card(slide, left, top, width, height, bg_color):
        shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
        shape.fill.solid()
        shape.fill.fore_color.rgb = bg_color
        shape.line.color.rgb = RGBColor(45, 45, 50) # Dark gray border
        shape.line.width = Pt(1)
        return shape

    # ==========================================
    # SLIDE 1 - CAPA
    # ==========================================
    slide = prs.slides.add_slide(blank_layout)
    set_dark_background(slide)
    
    # Center Logo Image
    slide.shapes.add_picture('logo.png', Inches(2.8), Inches(2.2), Inches(1.6), Inches(1.6))
    
    # Center Logo Text Box
    logo_text_box = slide.shapes.add_textbox(Inches(4.7), Inches(2.1), Inches(6.0), Inches(1.8))
    tf_logo = logo_text_box.text_frame
    tf_logo.word_wrap = True
    tf_logo.margin_left = Inches(0)
    tf_logo.margin_top = Inches(0)
    
    p_main = tf_logo.paragraphs[0]
    p_main.space_after = Pt(0)
    
    run1 = p_main.add_run()
    run1.text = "Suzano"
    run1.font.color.rgb = TEXT_WHITE
    run1.font.name = 'Arial'
    run1.font.size = Pt(68)
    run1.font.bold = True
    
    run2 = p_main.add_run()
    run2.text = " IT"
    run2.font.color.rgb = COLOR_ORANGE
    run2.font.name = 'Arial'
    run2.font.size = Pt(68)
    run2.font.bold = True
    
    # Slogan Centered Below
    slog_box = slide.shapes.add_textbox(Inches(1.666), Inches(4.3), Inches(10.0), Inches(1.0))
    tf_slog = slog_box.text_frame
    tf_slog.word_wrap = True
    p_slog = tf_slog.paragraphs[0]
    p_slog.alignment = PP_ALIGN.CENTER
    p_slog.text = "Tecnologia que move frotas."
    p_slog.font.name = 'Arial'
    p_slog.font.size = Pt(22)
    p_slog.font.color.rgb = TEXT_WHITE
    p_slog.font.bold = False
    
    add_footer(slide, 1)

    # ==========================================
    # SLIDE 2 - VISÃO GERAL DA PLATAFORMA
    # ==========================================
    slide = prs.slides.add_slide(blank_layout)
    set_dark_background(slide)
    add_header(slide, "Visão Geral da Plataforma", "Conceito e Integração")
    
    overview_intro = slide.shapes.add_textbox(Inches(0.8), Inches(1.7), Inches(11.7), Inches(0.8))
    tf_intro = overview_intro.text_frame
    p_intro = tf_intro.paragraphs[0]
    p_intro.text = "Uma solução robusta desenvolvida para unificar e modernizar todas as etapas operacionais do setor."
    p_intro.font.name = 'Arial'
    p_intro.font.size = Pt(18)
    p_intro.font.color.rgb = TEXT_MUTED
    
    # 3 Cards for the 3 main pillars
    pillars = [
        ("Multi-Tenant", "Plataforma SaaS multiempresa ideal para redes de oficinas, concessionárias e controle unificado de múltiplas filiais."),
        ("Inteligência Artificial", "CMH-AI aplicada para faturamento rápido de serviços, laudos inteligentes e pós-venda automotivo de alta precisão."),
        ("Gestão Unificada", "Controle centralizado de ordens de serviço, estoque de peças, módulos financeiros e gestão preventiva de frotas.")
    ]
    
    left_positions = [0.8, 4.8, 8.8]
    for i, (title, desc) in enumerate(pillars):
        draw_card(slide, Inches(left_positions[i]), Inches(2.6), Inches(3.733), Inches(3.8), COLOR_CARD_BG)
        box = slide.shapes.add_textbox(Inches(left_positions[i] + 0.1), Inches(2.8), Inches(3.533), Inches(3.4))
        tf_sol = box.text_frame
        tf_sol.word_wrap = True
        
        p_t = tf_sol.paragraphs[0]
        p_t.text = title
        p_t.font.name = 'Arial'
        p_t.font.size = Pt(18)
        p_t.font.bold = True
        p_t.font.color.rgb = COLOR_ORANGE
        p_t.space_after = Pt(12)
        
        p_d = tf_sol.add_paragraph()
        p_d.text = desc
        p_d.font.name = 'Arial'
        p_d.font.size = Pt(13)
        p_d.font.color.rgb = TEXT_MUTED
        
    add_footer(slide, 2)

    # ==========================================
    # SLIDE 3 - QUEM SOMOS
    # ==========================================
    slide = prs.slides.add_slide(blank_layout)
    set_dark_background(slide)
    add_header(slide, "Quem Somos", "Apresentação Institucional")
    
    # Left description box
    desc_box = slide.shapes.add_textbox(Inches(0.8), Inches(2.0), Inches(5.5), Inches(4.5))
    tf_desc = desc_box.text_frame
    tf_desc.word_wrap = True
    p_desc = tf_desc.paragraphs[0]
    p_desc.text = "A Suzano IT é uma empresa de tecnologia especializada no desenvolvimento de ecossistemas digitais para o setor automotivo.\n\nNosso propósito é reinventar a forma como oficinas mecânicas, auto centers e gestores de frota operam no dia a dia, eliminando controles manuais e proporcionando inteligência e redução de custos reais através de automação na nuvem."
    p_desc.font.name = 'Arial'
    p_desc.font.size = Pt(18)
    p_desc.font.color.rgb = TEXT_MUTED
    p_desc.space_after = Pt(14)
    
    # Right pillars box (cards layout)
    pillars = [
        ("Arquitetura 100% em Nuvem", "Acessível de qualquer lugar, dispensando servidores locais."),
        ("Multiempresa (Multi-Tenant)", "Ideal para redes de oficinas e controle de múltiplas filiais."),
        ("Inteligência Artificial", "CMH-AI aplicada à otimização e produtividade operacional.")
    ]
    
    top_offset = 2.0
    for title, text in pillars:
        draw_card(slide, Inches(6.8), Inches(top_offset), Inches(5.7), Inches(1.3), COLOR_CARD_BG)
        card_box = slide.shapes.add_textbox(Inches(7.0), Inches(top_offset + 0.1), Inches(5.3), Inches(1.1))
        tf_card = card_box.text_frame
        tf_card.word_wrap = True
        
        p_title = tf_card.paragraphs[0]
        p_title.text = title
        p_title.font.name = 'Arial'
        p_title.font.size = Pt(15)
        p_title.font.bold = True
        p_title.font.color.rgb = COLOR_SKY
        p_title.space_after = Pt(4)
        
        p_text = tf_card.add_paragraph()
        p_text.text = text
        p_text.font.name = 'Arial'
        p_text.font.size = Pt(12)
        p_text.font.color.rgb = TEXT_MUTED
        
        top_offset += 1.6
        
    add_footer(slide, 3)

    # ==========================================
    # SLIDE 3 - PROBLEMAS DO MERCADO
    # ==========================================
    slide = prs.slides.add_slide(blank_layout)
    set_dark_background(slide)
    add_header(slide, "Problemas do Mercado", "Gargalos do Setor")
    
    # Left Column: List of Pain Points
    left_box = slide.shapes.add_textbox(Inches(0.8), Inches(1.8), Inches(5.8), Inches(4.8))
    tf_left = left_box.text_frame
    tf_left.word_wrap = True
    
    p_head = tf_left.paragraphs[0]
    p_head.text = "Desafios em oficinas e gestão de frotas:"
    p_head.font.name = 'Arial'
    p_head.font.size = Pt(18)
    p_head.font.bold = True
    p_head.font.color.rgb = TEXT_WHITE
    p_head.space_after = Pt(14)
    
    pains = [
        "Controle manual de ordens de serviço e histórico lento.",
        "Falta de rastreabilidade na substituição e cotação de peças.",
        "Processos financeiros sem conciliação e com alta inadimplência.",
        "Falta de indicadores executivos de desempenho da operação.",
        "Desperdício de ativos em frotas por falta de preventiva.",
        "Sistemas antigos desconectados ou planilhas difíceis."
    ]
    for pain in pains:
        p_pain = tf_left.add_paragraph()
        p_pain.text = f"❌  {pain}"
        p_pain.font.name = 'Arial'
        p_pain.font.size = Pt(13)
        p_pain.font.color.rgb = TEXT_MUTED
        p_pain.space_after = Pt(10)
        
    # Right Column: Impacts Box (Warning Style)
    draw_card(slide, Inches(7.0), Inches(1.8), Inches(5.5), Inches(4.5), RGBColor(28, 18, 22))
    right_box = slide.shapes.add_textbox(Inches(7.3), Inches(2.1), Inches(4.9), Inches(3.9))
    tf_right = right_box.text_frame
    tf_right.word_wrap = True
    
    p_warn_head = tf_right.paragraphs[0]
    p_warn_head.text = "IMPACTOS FINANCEIROS DIRETOS"
    p_warn_head.font.name = 'Arial'
    p_warn_head.font.size = Pt(14)
    p_warn_head.font.bold = True
    p_warn_head.font.color.rgb = COLOR_RED
    p_warn_head.space_after = Pt(14)
    
    p_warn_body = tf_right.add_paragraph()
    p_warn_body.text = "A ineficiência operacional do controle manual gera desperdícios graves:\n\n•  Perda financeira de até 15% por vazamento de caixa e faturamento incorreto.\n•  Redução drástica de produtividade dos mecânicos e frotas paradas por quebra corretiva.\n•  Dificuldade gerencial por falta de dados para tomadas de decisão rápidas."
    p_warn_body.font.name = 'Arial'
    p_warn_body.font.size = Pt(15)
    p_warn_body.font.color.rgb = TEXT_MUTED
    p_warn_body.space_after = Pt(14)
    
    add_footer(slide, 4)

    # ==========================================
    # SLIDE 4 - NOSSA SOLUÇÃO
    # ==========================================
    slide = prs.slides.add_slide(blank_layout)
    set_dark_background(slide)
    add_header(slide, "Nossa Solução", "Ecossistema Suzano IT")
    
    sol_intro = slide.shapes.add_textbox(Inches(0.8), Inches(1.7), Inches(11.7), Inches(0.8))
    tf_intro = sol_intro.text_frame
    p_intro = tf_intro.paragraphs[0]
    p_intro.text = "Centralizamos toda a gestão operacional, de frotas e financeira em um único ecossistema em nuvem."
    p_intro.font.name = 'Arial'
    p_intro.font.size = Pt(18)
    p_intro.font.color.rgb = TEXT_MUTED
    
    # 4 Cards representing features
    solutions = [
        ("Gestão de Oficina", "Ordens de serviço digitais, controle completo de peças, cadastros de clientes e histórico por placa."),
        ("Gestão de Frotas", "Controle inteligente de pneus, consumo de combustível, motoristas e planos de manutenção preventiva."),
        ("Controle Financeiro", "Contas a pagar/receber integradas à operação, DRE estruturado e acompanhamento de fluxo de caixa."),
        ("Inteligência Artificial", "Modelagem preditiva CMH-AI para geração automática de NF-e e diagnósticos de desgaste.")
    ]
    
    left_positions = [0.8, 3.8, 6.8, 9.8]
    for i, (title, desc) in enumerate(solutions):
        draw_card(slide, Inches(left_positions[i]), Inches(2.6), Inches(2.733), Inches(3.8), COLOR_CARD_BG)
        box = slide.shapes.add_textbox(Inches(left_positions[i] + 0.1), Inches(2.8), Inches(2.533), Inches(3.4))
        tf_sol = box.text_frame
        tf_sol.word_wrap = True
        
        p_t = tf_sol.paragraphs[0]
        p_t.text = title
        p_t.font.name = 'Arial'
        p_t.font.size = Pt(18)
        p_t.font.bold = True
        p_t.font.color.rgb = COLOR_SKY
        p_t.space_after = Pt(12)
        
        p_d = tf_sol.add_paragraph()
        p_d.text = desc
        p_d.font.name = 'Arial'
        p_d.font.size = Pt(13)
        p_d.font.color.rgb = TEXT_MUTED
        
    add_footer(slide, 5)

    # ==========================================
    # SLIDE 5 - GESTÃO DE OFICINA
    # ==========================================
    slide = prs.slides.add_slide(blank_layout)
    set_dark_background(slide)
    add_header(slide, "Módulo: Gestão de Oficina", "Operação Integrada")
    
    # Left Column: Features
    left_box = slide.shapes.add_textbox(Inches(0.8), Inches(1.8), Inches(5.8), Inches(4.8))
    tf_left = left_box.text_frame
    tf_left.word_wrap = True
    
    p_head = tf_left.paragraphs[0]
    p_head.text = "Controle absoluto da oficina mecânica:"
    p_head.font.name = 'Arial'
    p_head.font.size = Pt(18)
    p_head.font.bold = True
    p_head.font.color.rgb = TEXT_WHITE
    p_head.space_after = Pt(14)
    
    features = [
        "Cadastro rápido de clientes e veículos integrados.",
        "Orçamentos comerciais digitais com envio em 1 clique.",
        "Ordens de serviço com controle de horas trabalhadas.",
        "Gestão de estoque de peças com reposição sugerida.",
        "Controle de comissões de mecânicos e consultores.",
        "Emissão de Notas Fiscais Eletrônicas de Serviço (NF-e)."
    ]
    for feat in features:
        p_feat = tf_left.add_paragraph()
        p_feat.text = f"✓  {feat}"
        p_feat.font.name = 'Arial'
        p_feat.font.size = Pt(14)
        p_feat.font.color.rgb = TEXT_MUTED
        p_feat.space_after = Pt(10)
        
    # Right Column: Results Card
    draw_card(slide, Inches(7.0), Inches(1.8), Inches(5.5), Inches(4.5), COLOR_CARD_BG)
    right_box = slide.shapes.add_textbox(Inches(7.3), Inches(2.1), Inches(4.9), Inches(3.9))
    tf_right = right_box.text_frame
    tf_right.word_wrap = True
    
    p_res_head = tf_right.paragraphs[0]
    p_res_head.text = "RESULTADOS OPERACIONAIS"
    p_res_head.font.name = 'Arial'
    p_res_head.font.size = Pt(14)
    p_res_head.font.bold = True
    p_res_head.font.color.rgb = COLOR_SKY
    p_res_head.space_after = Pt(18)
    
    p_res_body = tf_right.add_paragraph()
    p_res_body.text = "A digitalização de processos da Suzano IT gera impacto instantâneo:\n\n•  Organização completa do pátio mecânico.\n•  Melhoria média de 20% no faturamento bruto.\n•  Eliminação de erros de digitação e retrabalho.\n•  Aumento de confiança do cliente final com laudo online."
    p_res_body.font.name = 'Arial'
    p_res_body.font.size = Pt(15)
    p_res_body.font.color.rgb = TEXT_MUTED
    
    add_footer(slide, 6)

    # ==========================================
    # SLIDE 6 - GESTÃO DE FROTAS
    # ==========================================
    slide = prs.slides.add_slide(blank_layout)
    set_dark_background(slide)
    add_header(slide, "Módulo: Gestão de Frotas", "Controle de Ativos")
    
    # Left Column: Features
    left_box = slide.shapes.add_textbox(Inches(0.8), Inches(1.8), Inches(5.8), Inches(4.8))
    tf_left = left_box.text_frame
    tf_left.word_wrap = True
    
    p_head = tf_left.paragraphs[0]
    p_head.text = "Visão geral e preventiva da frota:"
    p_head.font.name = 'Arial'
    p_head.font.size = Pt(18)
    p_head.font.bold = True
    p_head.font.color.rgb = TEXT_WHITE
    p_head.space_after = Pt(14)
    
    features = [
        "Cadastro de veículos e subfrotas organizacionais.",
        "Gestão ativa de condutores e histórico de multas.",
        "Controle de abastecimento, média de consumo e CPK.",
        "Manutenção preventiva com alertas baseados em KM/Tempo.",
        "Controle de vida útil de pneus com histórico de rodagem.",
        "Rastreabilidade de custos e controle financeiro de ativos."
    ]
    for feat in features:
        p_feat = tf_left.add_paragraph()
        p_feat.text = f"✓  {feat}"
        p_feat.font.name = 'Arial'
        p_feat.font.size = Pt(14)
        p_feat.font.color.rgb = TEXT_MUTED
        p_feat.space_after = Pt(10)
        
    # Right Column: Results Card
    draw_card(slide, Inches(7.0), Inches(1.8), Inches(5.5), Inches(4.5), COLOR_CARD_BG)
    right_box = slide.shapes.add_textbox(Inches(7.3), Inches(2.1), Inches(4.9), Inches(3.9))
    tf_right = right_box.text_frame
    tf_right.word_wrap = True
    
    p_res_head = tf_right.paragraphs[0]
    p_res_head.text = "BENEFÍCIOS NA GESTÃO DE FROTA"
    p_res_head.font.name = 'Arial'
    p_res_head.font.size = Pt(14)
    p_res_head.font.bold = True
    p_res_head.font.color.rgb = COLOR_SKY
    p_res_head.space_after = Pt(18)
    
    p_res_body = tf_right.add_paragraph()
    p_res_body.text = "Economia em escala para empresas frotistas:\n\n•  Redução de custos operacionais diretos por manutenção planejada.\n•  Maior disponibilidade física dos veículos nas ruas.\n•  Prevenção de quebras graves e paradas indesejadas.\n•  Auditoria de consumo de combustível simplificada."
    p_res_body.font.name = 'Arial'
    p_res_body.font.size = Pt(15)
    p_res_body.font.color.rgb = TEXT_MUTED
    
    add_footer(slide, 7)

    # ==========================================
    # SLIDE 7 - INTELIGÊNCIA ARTIFICIAL
    # ==========================================
    slide = prs.slides.add_slide(blank_layout)
    set_dark_background(slide)
    add_header(slide, "Inteligência Artificial Integrada", "Tecnologia CMH-AI")
    
    left_box = slide.shapes.add_textbox(Inches(0.8), Inches(1.8), Inches(5.8), Inches(4.8))
    tf_left = left_box.text_frame
    tf_left.word_wrap = True
    
    p_head = tf_left.paragraphs[0]
    p_head.text = "Inteligência operacional no pós-venda:"
    p_head.font.name = 'Arial'
    p_head.font.size = Pt(18)
    p_head.font.bold = True
    p_head.font.color.rgb = TEXT_WHITE
    p_head.space_after = Pt(14)
    
    features = [
        "Geração de descrição de serviços para NF-e a partir de rascunhos informais.",
        "Sugestão preditiva de serviços adicionais baseada no modelo de uso do carro.",
        "Análise estatística para identificação preditiva de padrões de falhas mecânicas.",
        "Relatórios inteligentes para tomadas de decisão sem complexidade.",
        "Dashboards preditivos com alertas antecipados de manutenção crítica."
    ]
    for feat in features:
        p_feat = tf_left.add_paragraph()
        p_feat.text = f"✨  {feat}"
        p_feat.font.name = 'Arial'
        p_feat.font.size = Pt(14)
        p_feat.font.color.rgb = TEXT_MUTED
        p_feat.space_after = Pt(10)
        
    # Right Column: AI Simulation Example
    draw_card(slide, Inches(7.0), Inches(1.8), Inches(5.5), Inches(4.5), COLOR_CARD_BG)
    right_box = slide.shapes.add_textbox(Inches(7.3), Inches(2.1), Inches(4.9), Inches(3.9))
    tf_right = right_box.text_frame
    tf_right.word_wrap = True
    
    p_res_head = tf_right.paragraphs[0]
    p_res_head.text = "EXEMPLO DE USO PRÁTICO (CMH-AI)"
    p_res_head.font.name = 'Arial'
    p_res_head.font.size = Pt(14)
    p_res_head.font.bold = True
    p_res_head.font.color.rgb = COLOR_SKY
    p_res_head.space_after = Pt(10)
    
    p_res_body = tf_right.add_paragraph()
    p_res_body.text = "Rascunho de Entrada do Mecânico:\n'troca oleo velas filtro ar Corolla'\n\nRetorno da IA de Faturamento:\n'Prestação de serviços mecânicos englobando diagnóstico computadorizado de injeção, substituição preventiva do filtro de ar e velas de ignição originais, com revisão detalhada do plano de lubrificação do motor.'\n\nBenefício: Economia de tempo administrativa de recepção."
    p_res_body.font.name = 'Arial'
    p_res_body.font.size = Pt(13)
    p_res_body.font.color.rgb = TEXT_MUTED
    
    add_footer(slide, 8)

    # ==========================================
    # SLIDE 8 - DIFERENCIAIS
    # ==========================================
    slide = prs.slides.add_slide(blank_layout)
    set_dark_background(slide)
    add_header(slide, "Diferenciais Competitivos", "Por que a Suzano IT?")
    
    diffs = [
        ("SaaS Multiempresa", "Isolamento lógico absoluto e suporte completo para franquias e filiais."),
        ("Tabela FIPE Nativa", "Valores atualizados de mercado integrados diretamente às fichas do veículo."),
        ("ReceitaWS Integrada", "Preenchimento imediato de cadastros de pessoas jurídicas usando apenas CNPJ."),
        ("Faturamento NF-e", "Emissão tributária ágil em conformidade com as legislações vigentes."),
        ("Segurança & LGPD", "Dados protegidos com criptografia de alto nível e políticas rígidas de acesso."),
        ("Cloud Elasticidade", "Hospedagem elástica AWS e atualizações automáticas sem necessidade de intervenção.")
    ]
    
    x_positions = [0.8, 4.8, 8.8]
    y_positions = [2.0, 4.4]
    
    for idx, (title, desc) in enumerate(diffs):
        col = idx % 3
        row = idx // 3
        left = x_positions[col]
        top = y_positions[row]
        
        draw_card(slide, Inches(left), Inches(top), Inches(3.733), Inches(2.0), COLOR_CARD_BG)
        box = slide.shapes.add_textbox(Inches(left + 0.1), Inches(top + 0.1), Inches(3.533), Inches(1.8))
        tf_card = box.text_frame
        tf_card.word_wrap = True
        
        p_t = tf_card.paragraphs[0]
        p_t.text = title
        p_t.font.name = 'Arial'
        p_t.font.size = Pt(16)
        p_t.font.bold = True
        p_t.font.color.rgb = COLOR_SKY
        p_t.space_after = Pt(8)
        
        p_d = tf_card.add_paragraph()
        p_d.text = desc
        p_d.font.name = 'Arial'
        p_d.font.size = Pt(12)
        p_d.font.color.rgb = TEXT_MUTED
        
    add_footer(slide, 9)

    # ==========================================
    # SLIDE 9 - MODELO DE NEGÓCIO
    # ==========================================
    slide = prs.slides.add_slide(blank_layout)
    set_dark_background(slide)
    add_header(slide, "Modelo de Negócio", "Monetização SaaS")
    
    intro_box = slide.shapes.add_textbox(Inches(0.8), Inches(1.6), Inches(11.7), Inches(0.6))
    tf_intro = intro_box.text_frame
    p_intro = tf_intro.paragraphs[0]
    p_intro.text = "Assinaturas mensais recorrentes (MRR) estruturadas para crescer conforme o uso e porte do cliente."
    p_intro.font.name = 'Arial'
    p_intro.font.size = Pt(16)
    p_intro.font.color.rgb = TEXT_MUTED
    
    plans = [
        ("STARTER", "R$ 199 / mês", "Pequenas Oficinas", [
            "Funcionalidades essenciais",
            "Gestão de OS e Cadastros",
            "Orçamentos comerciais",
            "Suporte digital básico"
        ]),
        ("PROFESSIONAL", "R$ 499 / mês", "Empresas em Crescimento", [
            "Gestão de Oficina completa",
            "Módulo de Frotas básico",
            "Controle financeiro avançado",
            "Emissão integrada de NF-e"
        ]),
        ("ENTERPRISE", "Sob Consulta", "Grandes Operações", [
            "Frotistas e transportadoras",
            "Integrações de API personalizadas",
            "Arquitetura sob demanda",
            "Suporte e SLA dedicado"
        ])
    ]
    
    left_positions = [0.8, 4.8, 8.8]
    for i, (name, price, target, features) in enumerate(plans):
        bg = COLOR_CARD_BG if name != "PROFESSIONAL" else RGBColor(18, 30, 49)
        draw_card(slide, Inches(left_positions[i]), Inches(2.3), Inches(3.733), Inches(4.3), bg)
        
        box = slide.shapes.add_textbox(Inches(left_positions[i] + 0.1), Inches(2.5), Inches(3.533), Inches(3.9))
        tf_plan = box.text_frame
        tf_plan.word_wrap = True
        
        p_n = tf_plan.paragraphs[0]
        p_n.text = name
        p_n.font.name = 'Arial'
        p_n.font.size = Pt(18)
        p_n.font.bold = True
        p_n.font.color.rgb = COLOR_SKY if name != "PROFESSIONAL" else COLOR_INDIGO
        p_n.space_after = Pt(4)
        
        p_p = tf_plan.add_paragraph()
        p_p.text = price
        p_p.font.name = 'Arial'
        p_p.font.size = Pt(22)
        p_p.font.bold = True
        p_p.font.color.rgb = TEXT_WHITE
        p_p.space_after = Pt(2)
        
        p_t = tf_plan.add_paragraph()
        p_t.text = target
        p_t.font.name = 'Arial'
        p_t.font.size = Pt(12)
        p_t.font.color.rgb = TEXT_MUTED
        p_t.space_after = Pt(18)
        
        for feat in features:
            p_f = tf_plan.add_paragraph()
            p_f.text = f"•  {feat}"
            p_f.font.name = 'Arial'
            p_f.font.size = Pt(12)
            p_f.font.color.rgb = TEXT_MUTED
            p_f.space_after = Pt(6)
            
    add_footer(slide, 10)

    # ==========================================
    # SLIDE 10 - MERCADO
    # ==========================================
    slide = prs.slides.add_slide(blank_layout)
    set_dark_background(slide)
    add_header(slide, "Mercado Alvo", "Público e Potencial")
    
    left_box = slide.shapes.add_textbox(Inches(0.8), Inches(1.8), Inches(5.8), Inches(4.8))
    tf_left = left_box.text_frame
    tf_left.word_wrap = True
    
    p_head = tf_left.paragraphs[0]
    p_head.text = "Quem são nossos clientes:"
    p_head.font.name = 'Arial'
    p_head.font.size = Pt(18)
    p_head.font.bold = True
    p_head.font.color.rgb = TEXT_WHITE
    p_head.space_after = Pt(14)
    
    targets = [
        "Oficinas mecânicas e auto centers independentes.",
        "Transportadoras, frotas corporativas e logísticas.",
        "Locadoras de veículos e cooperativas de táxi/aplicativos.",
        "Concessionárias de veículos e redes de autopeças.",
        "Prefeituras municipais e órgãos governamentais."
    ]
    for target in targets:
        p_t = tf_left.add_paragraph()
        p_t.text = f"•  {target}"
        p_t.font.name = 'Arial'
        p_t.font.size = Pt(14)
        p_t.font.color.rgb = TEXT_MUTED
        p_t.space_after = Pt(12)
        
    # Right Column
    draw_card(slide, Inches(7.0), Inches(1.8), Inches(5.5), Inches(4.5), COLOR_CARD_BG)
    right_box = slide.shapes.add_textbox(Inches(7.3), Inches(2.1), Inches(4.9), Inches(3.9))
    tf_right = right_box.text_frame
    tf_right.word_wrap = True
    
    p_res_head = tf_right.paragraphs[0]
    p_res_head.text = "POTENCIAL DE CRESCIMENTO"
    p_res_head.font.name = 'Arial'
    p_res_head.font.size = Pt(14)
    p_res_head.font.bold = True
    p_res_head.font.color.rgb = COLOR_SKY
    p_res_head.space_after = Pt(18)
    
    p_res_body = tf_right.add_paragraph()
    p_res_body.text = "O mercado brasileiro conta com mais de 100 mil oficinas de pós-venda cadastradas e um crescimento acelerado no número de empresas optando por gerenciar frotas próprias para reduzir custos de distribuição.\n\nA digitalização e a inteligência operacional são as maiores prioridades para o setor nos próximos anos."
    p_res_body.font.name = 'Arial'
    p_res_body.font.size = Pt(15)
    p_res_body.font.color.rgb = TEXT_MUTED
    
    add_footer(slide, 11)

    # ==========================================
    # SLIDE 11 - ARQUITETURA
    # ==========================================
    slide = prs.slides.add_slide(blank_layout)
    set_dark_background(slide)
    add_header(slide, "Arquitetura Tecnológica", "Tecnologia e Performance")
    
    # 4 Components Side by Side
    techs = [
        ("Frontend", "React\nTypeScript\nTailwind CSS\n\n•  Interface SPA responsiva e veloz."),
        ("Backend", "Kotlin\nSpring Boot\nJava Virtual Machine\n\n•  Segurança robusta e performance corporativa."),
        ("Banco de Dados", "PostgreSQL\n\n\n•  Persistência relacional estável com alta auditabilidade."),
        ("Segurança & Cloud", "JWT Autenticação\nDocker Containers\nAWS Cloud\n\n•  Segurança corporativa aderente à LGPD.")
    ]
    
    left_positions = [0.8, 3.8, 6.8, 9.8]
    for i, (name, details) in enumerate(techs):
        draw_card(slide, Inches(left_positions[i]), Inches(2.0), Inches(2.733), Inches(4.3), COLOR_CARD_BG)
        box = slide.shapes.add_textbox(Inches(left_positions[i] + 0.1), Inches(2.2), Inches(2.533), Inches(3.9))
        tf_tech = box.text_frame
        tf_tech.word_wrap = True
        
        p_name = tf_tech.paragraphs[0]
        p_name.text = name
        p_name.font.name = 'Arial'
        p_name.font.size = Pt(18)
        p_name.font.bold = True
        p_name.font.color.rgb = COLOR_SKY
        p_name.space_after = Pt(14)
        
        p_det = tf_tech.add_paragraph()
        p_det.text = details
        p_det.font.name = 'Arial'
        p_det.font.size = Pt(12)
        p_det.font.color.rgb = TEXT_MUTED
        
    add_footer(slide, 12)

    # ==========================================
    # SLIDE 12 - ROADMAP
    # ==========================================
    slide = prs.slides.add_slide(blank_layout)
    set_dark_background(slide)
    add_header(slide, "Roadmap de Evolução", "Cronograma de Lançamento")
    
    # 7 Columns represent phases
    phases = [
        ("Fase 1", "Gestão Oficina", "✓ Concluído"),
        ("Fase 2", "Gestão Frota", "✓ Concluído"),
        ("Fase 3", "Integração NF-e", "✓ Concluído"),
        ("Fase 4", "Inteligência Artif.", "✓ Concluído"),
        ("Fase 5", "Aplicativo Mobile", "Próximo"),
        ("Fase 6", "Marketplace", "Futuro"),
        ("Fase 7", "BI & Analytics", "Futuro")
    ]
    
    left_positions = [0.8, 2.5, 4.2, 5.9, 7.6, 9.3, 11.0]
    for i, (name, title, status) in enumerate(phases):
        bg = RGBColor(16, 27, 23) if "Concluído" in status else (RGBColor(16, 23, 40) if "Próximo" in status else COLOR_CARD_BG)
        draw_card(slide, Inches(left_positions[i]), Inches(2.3), Inches(1.533), Inches(3.8), bg)
        
        box = slide.shapes.add_textbox(Inches(left_positions[i] + 0.05), Inches(2.5), Inches(1.433), Inches(3.4))
        tf_p = box.text_frame
        tf_p.word_wrap = True
        
        p_n = tf_p.paragraphs[0]
        p_n.text = name
        p_n.font.name = 'Arial'
        p_n.font.size = Pt(16)
        p_n.font.bold = True
        p_n.font.color.rgb = TEXT_WHITE
        p_n.space_after = Pt(14)
        
        p_t = tf_p.add_paragraph()
        p_t.text = title
        p_t.font.name = 'Arial'
        p_t.font.size = Pt(12)
        p_t.font.bold = True
        p_t.font.color.rgb = COLOR_SKY
        p_t.space_after = Pt(18)
        
        p_s = tf_p.add_paragraph()
        p_s.text = status
        p_s.font.name = 'Arial'
        p_s.font.size = Pt(11)
        p_s.font.color.rgb = RGBColor(34, 197, 94) if "Concluído" in status else (COLOR_SKY if "Próximo" in status else TEXT_MUTED)
        
    add_footer(slide, 13)

    # ==========================================
    # SLIDE 13 - BENEFÍCIOS PARA O CLIENTE
    # ==========================================
    slide = prs.slides.add_slide(blank_layout)
    set_dark_background(slide)
    add_header(slide, "Benefícios para o Cliente", "Resultados Finais")
    
    # 4 Cards side by side
    benefits = [
        ("Redução de Custos", "Evita quebras inesperadas através de cronogramas rígidos de manutenção preventiva na frota."),
        ("Organização Geral", "Fluxo operacional padronizado desde a abertura da OS até o faturamento e envio de laudos."),
        ("Decisões com Dados", "Indicadores e relatórios consolidados em dashboards executivos em tempo real."),
        ("Produtividade Elevada", "Eliminação de controles em papéis e digitação manual por meio de automação e Inteligência Artificial.")
    ]
    
    left_positions = [0.8, 3.8, 6.8, 9.8]
    for i, (title, desc) in enumerate(benefits):
        draw_card(slide, Inches(left_positions[i]), Inches(2.2), Inches(2.733), Inches(3.9), COLOR_CARD_BG)
        box = slide.shapes.add_textbox(Inches(left_positions[i] + 0.1), Inches(2.4), Inches(2.533), Inches(3.5))
        tf_b = box.text_frame
        tf_b.word_wrap = True
        
        p_t = tf_b.paragraphs[0]
        p_t.text = title
        p_t.font.name = 'Arial'
        p_t.font.size = Pt(18)
        p_t.font.bold = True
        p_t.font.color.rgb = COLOR_SKY
        p_t.space_after = Pt(12)
        
        p_d = tf_b.add_paragraph()
        p_d.text = desc
        p_d.font.name = 'Arial'
        p_d.font.size = Pt(13)
        p_d.font.color.rgb = TEXT_MUTED
        
    add_footer(slide, 14)

    # ==========================================
    # SLIDE 15 - ENCERRAMENTO
    # ==========================================
    slide = prs.slides.add_slide(blank_layout)
    set_dark_background(slide)
    
    draw_card(slide, Inches(2.0), Inches(1.5), Inches(9.333), Inches(4.8), COLOR_CARD_BG)
    
    # Center Logo Image inside the card
    slide.shapes.add_picture('logo.png', Inches(4.3), Inches(2.0), Inches(1.2), Inches(1.2))
    
    # Center Logo Text Box inside the card
    logo_text_box = slide.shapes.add_textbox(Inches(5.7), Inches(1.95), Inches(4.0), Inches(1.3))
    tf_logo = logo_text_box.text_frame
    tf_logo.word_wrap = True
    tf_logo.margin_left = Inches(0)
    tf_logo.margin_top = Inches(0)
    p_main = tf_logo.paragraphs[0]
    p_main.space_after = Pt(0)
    
    run1 = p_main.add_run()
    run1.text = "Suzano"
    run1.font.color.rgb = TEXT_WHITE
    run1.font.name = 'Arial'
    run1.font.size = Pt(48)
    run1.font.bold = True
    
    run2 = p_main.add_run()
    run2.text = " IT"
    run2.font.color.rgb = COLOR_ORANGE
    run2.font.name = 'Arial'
    run2.font.size = Pt(48)
    run2.font.bold = True
    
    # Slogan Centered Below Logo inside the card
    slog_box = slide.shapes.add_textbox(Inches(2.5), Inches(3.4), Inches(8.333), Inches(0.6))
    tf_slog = slog_box.text_frame
    tf_slog.word_wrap = True
    p_slog = tf_slog.paragraphs[0]
    p_slog.alignment = PP_ALIGN.CENTER
    p_slog.text = "Tecnologia que move frotas."
    p_slog.font.name = 'Arial'
    p_slog.font.size = Pt(18)
    p_slog.font.color.rgb = TEXT_WHITE
    p_slog.font.bold = False
    
    # Contact Info centered at bottom of the card
    contact_box = slide.shapes.add_textbox(Inches(2.5), Inches(4.1), Inches(8.333), Inches(1.2))
    tf_contact = contact_box.text_frame
    tf_contact.word_wrap = True
    
    p_contact = tf_contact.paragraphs[0]
    p_contact.alignment = PP_ALIGN.CENTER
    p_contact.text = "Website: www.suzanoit.com.br   |   E-mail: contato@suzanoit.com.br"
    p_contact.font.name = 'Arial'
    p_contact.font.size = Pt(13)
    p_contact.font.color.rgb = TEXT_MUTED
    p_contact.space_after = Pt(14)
    
    p_quote = tf_contact.add_paragraph()
    p_quote.alignment = PP_ALIGN.CENTER
    p_quote.text = '"Mais controle. Mais produtividade. Mais resultado."'
    p_quote.font.name = 'Arial'
    p_quote.font.size = Pt(14)
    p_quote.font.bold = True
    p_quote.font.color.rgb = COLOR_ORANGE
    
    add_footer(slide, 15)

    # Save presentation
    filename = "Apresentacao_Suzano_IT.pptx"
    prs.save(filename)
    print(f"Presentation saved successfully as {filename}")

if __name__ == "__main__":
    create_presentation()

